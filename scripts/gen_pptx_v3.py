"""
Generate v3 PowerPoint: Add GitHub Actions content slides and comparison slide.
Replaces slides 14-17 with content-rich GitHub Actions slides and adds a comparison slide.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import copy
from lxml import etree
import os

SRC = r'C:\Users\darioa\Microsoft\Microsoft (Internal)-12 - Projects - PresentazioneDevops - Documents\AI Lab for Startup - GitHub + Azure - CI CD - v2a.pptx'
DST = r'C:\Users\darioa\Microsoft\Microsoft (Internal)-12 - Projects - PresentazioneDevops - Documents\AI Lab for Startup - GitHub + Azure - CI CD - v3.pptx'

prs = Presentation(SRC)
blank_layout = prs.slide_layouts[6]  # Blank layout

# ═══════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════

def add_left_bar(slide, width=164592):
    """Add the dark-blue left bar consistent with existing slides."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1B, 0x1B, 0x1B)
    shape.line.fill.background()
    shape.name = "Rectangle 8"

def add_title(slide, text, left=430225, top=178385, width=10972800, height=646331, size=Pt(36), bold=True, color=None):
    """Add a slide title text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    return txBox

def add_title_underline(slide, left=512670, top=818464, width=10091420, height=45719):
    """Add the horizontal line under the title."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x00, 0x78, 0xD4)  # Microsoft blue
    shape.line.fill.background()

def add_section_header(slide, text, left, top, width=2028440, height=430887, size=Pt(22), bold=True):
    """Add a section header like 'Key Capabilities'."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    return txBox

def add_vertical_divider(slide, left=5989320, top=1097280, width=18288, height=5486400):
    """Add vertical divider between left and right columns."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    shape.line.fill.background()

def add_bullet_item(slide, title_text, desc1, desc2, left, top, 
                    bullet_shape=MSO_SHAPE.OVAL, bullet_color=RGBColor(0x00, 0x78, 0xD4),
                    title_size=Pt(13), desc_size=Pt(10)):
    """Add a bullet point with title and two description lines."""
    # Bullet shape
    bullet_left = left
    bullet_top = top + 91440  # offset for visual alignment
    bshape = slide.shapes.add_shape(bullet_shape, bullet_left, bullet_top, 109728, 109728)
    bshape.fill.solid()
    bshape.fill.fore_color.rgb = bullet_color
    bshape.line.fill.background()
    
    # Text box
    text_left = left + 182880
    txBox = slide.shapes.add_textbox(text_left, top, 5029200, 640080)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Title line
    p = tf.paragraphs[0]
    p.text = title_text
    run = p.runs[0]
    run.font.size = title_size
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    
    # Description line 1
    p2 = tf.add_paragraph()
    p2.text = desc1
    run2 = p2.runs[0]
    run2.font.size = desc_size
    run2.font.bold = False
    run2.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
    
    # Description line 2
    p3 = tf.add_paragraph()
    p3.text = desc2
    run3 = p3.runs[0]
    run3.font.size = desc_size
    run3.font.bold = False
    run3.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

def add_code_block(slide, code_text, left, top, width, height, font_size=Pt(9)):
    """Add a code block with monospace font and dark background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)
    
    lines = code_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = font_size
            run.font.name = 'Consolas'
            run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)

def add_table_slide(slide, headers, rows, left, top, width, col_widths=None):
    """Add a table to a slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    row_height = Emu(320040)
    table_height = row_height * num_rows
    
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, table_height)
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x00, 0x78, 0xD4)
    
    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

def add_subtitle_text(slide, text, left, top, width=10515600, height=369332, size=Pt(14), bold=False, color=None):
    """Add a subtitle / descriptive text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
    return txBox

# ═══════════════════════════════════════════
# DELETE EXISTING SLIDES 14-17 (0-indexed: 13-16) AND INSERT NEW ONES
# ═══════════════════════════════════════════

# We'll delete slides 14, 15, 16, 17 (1-indexed) and insert our new slides in their place.
# python-pptx doesn't have a direct "delete slide" API, so we manipulate the XML.

slide_list = prs.slides._sldIdLst
slide_ids = list(slide_list)

# Slides to remove: indices 13, 14, 15, 16 (0-based) = slides 14, 15, 16, 17
# We need to remove them from the presentation relationship and slide list
slides_to_remove = [13, 14, 15, 16]  # 0-indexed

# Get the rIds and remove slides
removed_rids = []
for idx in sorted(slides_to_remove, reverse=True):
    sldId = slide_ids[idx]
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    removed_rids.append(rId)
    slide_list.remove(sldId)

# Remove the relationships and parts
for rId in removed_rids:
    rel = prs.part.rels[rId]
    part = rel.target_part
    # Remove the part from the package
    prs.part.drop_rel(rId)

# ═══════════════════════════════════════════
# NOW ADD NEW SLIDES (they will be appended at the end, then we reorder)
# ═══════════════════════════════════════════

# The slides are appended to the end. After adding all, we'll reorder them.
# Current slide count after removal: 24 (was 28, removed 4)
# New slides should go at position 14 (1-indexed), i.e., after slide 13

new_slides = []

# ═══════════════════════════════════════════
# SLIDE A: GitHub Actions — At a Glance (Key Capabilities + Key Limitations)
# Mirror of slide 11 (Azure DevOps At a Glance)
# ═══════════════════════════════════════════
slide_a = prs.slides.add_slide(blank_layout)
new_slides.append(slide_a)

add_left_bar(slide_a)
add_title(slide_a, "GitHub Actions — At a Glance", left=429768, top=182880)
add_title_underline(slide_a, left=522706, top=797177, width=7262152, height=36576)

# Vertical divider
add_vertical_divider(slide_a)

# Left column: Key Capabilities
add_section_header(slide_a, "Key Capabilities", left=548640, top=1051560)

cap_items = [
    ("Event-driven CI/CD",
     "25+ trigger types: push, PR, schedule, dispatch, release, etc.",
     "Native GitHub integration — status checks, PR comments, deployments"),
    ("Reusable workflows + composite actions",
     "Modular: reusable workflows (jobs) + composite actions (steps)",
     "Cross-repo workflow calls from centralized template repos"),
    ("Matrix strategy",
     "Generate N parallel jobs from OS × version × config combinations",
     "Dynamic matrix via fromJSON() for runtime-computed test grids"),
    ("OIDC — secretless cloud auth",
     "Federated credentials for Azure, AWS, GCP — no stored secrets",
     "Token claims scoped by repo, branch, environment"),
    ("20,000+ marketplace actions",
     "Rich ecosystem: setup-dotnet, azure/login, docker/build-push, etc.",
     "Composite, JavaScript, or Docker container action types"),
    ("Self-hosted runners (free & unlimited)",
     "No per-agent charge, no minute billing on your infrastructure",
     "Actions Runner Controller (ARC) for Kubernetes auto-scaling"),
    ("Built-in supply chain security",
     "Dependabot, CodeQL SAST, secret scanning, push protection",
     "SLSA provenance attestations + dependency review in PRs"),
    ("GitHub Pages & Packages",
     "Built-in artifact hosting (npm, NuGet, Docker, Maven)",
     "GitHub Pages deployment from workflows with environment gates"),
]

y_start = 1490472
y_step = 640080
for i, (title, d1, d2) in enumerate(cap_items):
    add_bullet_item(slide_a, title, d1, d2,
                    left=594360, top=y_start + i * y_step,
                    bullet_shape=MSO_SHAPE.OVAL,
                    bullet_color=RGBColor(0x00, 0x78, 0xD4))

# Right column: Key Limitations
add_section_header(slide_a, "Key Limitations", left=6263640, top=1051560)

lim_items = [
    ("Free tier: 2,000 min/month (private repos)",
     "Linux ×1, Windows ×2, macOS ×10 minute multipliers",
     "Public repos: unlimited free minutes"),
    ("No stage concept",
     "Flat job graph only — stages simulated via needs: dependencies",
     "No visual stage grouping in the GitHub UI"),
    ("Reusable workflow depth: max 4 levels",
     "Cannot nest reusable workflows deeper than 4 calls",
     "Reusable workflows produce jobs, not individual steps"),
    ("Limited input types",
     "Only string, boolean, number — no object/array parameters",
     "Workaround: JSON-serialize via fromJSON()"),
    ("Job timeout: 6 hrs hosted, 35 days self-hosted",
     "GitHub-hosted max 360 min per job — no override available",
     "Self-hosted: up to 35 days per job (effectively unlimited)"),
    ("Runner environments",
     "GitHub-hosted: Linux (x64/ARM64), Windows (x64), macOS (x64/ARM64)",
     "Larger runners (4-64 vCPU, GPU) require Team/Enterprise plan"),
    ("Scheduled workflows: default branch only",
     "Cron triggers only run on the default branch",
     "Disabled after 60 days of repo inactivity"),
    ("No enforced template inheritance",
     "No extends keyword — cannot require all repos use approved templates",
     "Org-level action restrictions (allowlist) are the closest control"),
]

for i, (title, d1, d2) in enumerate(lim_items):
    add_bullet_item(slide_a, title, d1, d2,
                    left=6309359, top=y_start + i * y_step,
                    bullet_shape=MSO_SHAPE.DIAMOND,
                    bullet_color=RGBColor(0xD8, 0x3B, 0x01))

# ═══════════════════════════════════════════
# SLIDE B: GitHub Actions — Workflow Structure (with code example)
# ═══════════════════════════════════════════
slide_b = prs.slides.add_slide(blank_layout)
new_slides.append(slide_b)

add_left_bar(slide_b)
add_title(slide_b, "Structure of a GitHub Actions Workflow", left=430225, top=178385)
add_title_underline(slide_b, left=512670, top=818464)

# Left: workflow structure code
code_left = """name: CI/CD Pipeline
on:
  push:
    branches: [main]
    paths: ['src/**']
  pull_request:
    branches: [main]
  workflow_dispatch:
    inputs:
      environment:
        type: choice
        options: [dev, staging, prod]

permissions:
  contents: read
  id-token: write

concurrency:
  group: ${{ github.workflow }}-${{ github.ref }}
  cancel-in-progress: true

env:
  DOTNET_VERSION: '9.0.x'"""

add_code_block(slide_b, code_left, 
               left=Emu(300000), top=Emu(950000), 
               width=Emu(5400000), height=Emu(5600000),
               font_size=Pt(8))

# Right: jobs structure code
code_right = """jobs:
  build:
    runs-on: ubuntu-latest
    outputs:
      version: ${{ steps.ver.outputs.version }}
    steps:
      - uses: actions/checkout@v4
      - uses: actions/setup-dotnet@v4
        with:
          dotnet-version: ${{ env.DOTNET_VERSION }}
      - run: dotnet build -c Release
      - id: ver
        run: echo "version=1.0.${{ github.run_number }}" \\
              >> $GITHUB_OUTPUT
      - uses: actions/upload-artifact@v4
        with:
          name: build-output
          path: ./publish/

  deploy:
    needs: build
    runs-on: ubuntu-latest
    environment:
      name: production
      url: https://myapp.azurewebsites.net
    steps:
      - uses: azure/login@v2
        with:
          client-id: ${{ vars.AZURE_CLIENT_ID }}
          tenant-id: ${{ vars.AZURE_TENANT_ID }}
          subscription-id: ${{ vars.AZURE_SUBSCRIPTION_ID }}
      - uses: azure/webapps-deploy@v3
        with:
          app-name: my-app"""

add_code_block(slide_b, code_right,
               left=Emu(5900000), top=Emu(950000),
               width=Emu(5900000), height=Emu(5600000),
               font_size=Pt(8))

# Annotations
annotations = [
    (Emu(300000), Emu(900000), "Triggers & permissions"),
    (Emu(5900000), Emu(900000), "Jobs, steps & deployment"),
]
for ax, ay, atext in annotations:
    txb = slide_b.shapes.add_textbox(ax, ay, Emu(3000000), Emu(200000))
    p = txb.text_frame.paragraphs[0]
    p.text = atext
    run = p.runs[0]
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x78, 0xD4)

# ═══════════════════════════════════════════
# SLIDE C: GitHub Actions — Reusable Workflows & Advanced Patterns
# ═══════════════════════════════════════════
slide_c = prs.slides.add_slide(blank_layout)
new_slides.append(slide_c)

add_left_bar(slide_c)
add_title(slide_c, "GitHub Actions — Reusable Workflows & Patterns", left=430225, top=178385)
add_title_underline(slide_c, left=512670, top=818464)

# Left: Reusable workflow definition
code_reusable = """# .github/workflows/deploy.yml (reusable)
name: Deploy to Azure
on:
  workflow_call:
    inputs:
      environment:
        required: true
        type: string
      app-name:
        required: true
        type: string
    outputs:
      url:
        value: ${{ jobs.deploy.outputs.url }}

jobs:
  deploy:
    runs-on: ubuntu-latest
    environment:
      name: ${{ inputs.environment }}
    outputs:
      url: ${{ steps.deploy.outputs.webapp-url }}
    steps:
      - uses: azure/login@v2
        with:
          client-id: ${{ vars.AZURE_CLIENT_ID }}
          tenant-id: ${{ vars.AZURE_TENANT_ID }}
          subscription-id: ${{ vars.AZURE_SUBSCRIPTION_ID }}
      - uses: azure/webapps-deploy@v3
        id: deploy
        with:
          app-name: ${{ inputs.app-name }}"""

add_code_block(slide_c, code_reusable,
               left=Emu(300000), top=Emu(950000),
               width=Emu(5400000), height=Emu(4800000),
               font_size=Pt(8))

# Right: Caller + matrix + composite
code_caller = """# Caller workflow — uses reusable workflow
jobs:
  deploy-staging:
    uses: ./.github/workflows/deploy.yml
    with:
      environment: staging
      app-name: my-app-staging
    secrets: inherit

  deploy-prod:
    needs: deploy-staging
    uses: ./.github/workflows/deploy.yml
    with:
      environment: production
      app-name: my-app-prod
    secrets: inherit

# Matrix strategy example
  test:
    strategy:
      fail-fast: false
      matrix:
        os: [ubuntu-latest, windows-latest]
        dotnet: ['8.0.x', '9.0.x']
    runs-on: ${{ matrix.os }}
    steps:
      - uses: actions/setup-dotnet@v4
        with:
          dotnet-version: ${{ matrix.dotnet }}
      - run: dotnet test"""

add_code_block(slide_c, code_caller,
               left=Emu(5900000), top=Emu(950000),
               width=Emu(5900000), height=Emu(4800000),
               font_size=Pt(8))

# Annotations
ann2 = [
    (Emu(300000), Emu(900000), "Reusable workflow definition (callee)"),
    (Emu(5900000), Emu(900000), "Caller workflow + matrix strategy"),
]
for ax, ay, atext in ann2:
    txb = slide_c.shapes.add_textbox(ax, ay, Emu(4000000), Emu(200000))
    p = txb.text_frame.paragraphs[0]
    p.text = atext
    run = p.runs[0]
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x78, 0xD4)

# Key patterns summary
patterns_text = [
    "▸ Reusable workflows (workflow_call) = job-level reuse across repos",
    "▸ Composite actions = step-level reuse (local or cross-repo)",
    "▸ Matrix strategy = N parallel jobs from OS × version × config",
    "▸ Dynamic matrix = fromJSON() for runtime-computed job grids",
    "▸ OIDC auth = secretless Azure/AWS/GCP via federated credentials",
    "▸ secrets: inherit = pass all caller secrets without listing each one",
]
txBox = slide_c.shapes.add_textbox(Emu(300000), Emu(5900000), Emu(11500000), Emu(800000))
tf = txBox.text_frame
tf.word_wrap = True
for i, pt_text in enumerate(patterns_text):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = pt_text
    run = p.runs[0]
    run.font.size = Pt(9)
    run.font.bold = False
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

# ═══════════════════════════════════════════
# SLIDE D: GitHub Actions — Security & Environments
# ═══════════════════════════════════════════
slide_d = prs.slides.add_slide(blank_layout)
new_slides.append(slide_d)

add_left_bar(slide_d)
add_title(slide_d, "GitHub Actions — Security & Environments", left=430225, top=178385)
add_title_underline(slide_d, left=512670, top=818464)

# Left column: Security features
sec_items = [
    ("OIDC Federated Credentials",
     "Secretless auth to Azure, AWS, GCP — no stored credentials",
     "Token scoped by repo + branch + environment claims"),
    ("GITHUB_TOKEN (auto-provisioned)",
     "Per-workflow token with configurable permission scopes",
     "Least-privilege: contents:read, id-token:write, etc."),
    ("Secrets at 3 levels",
     "Organization → Repository → Environment scoped secrets",
     "Masked in logs, blocked from fork PRs"),
    ("Environment protection rules",
     "Required reviewers, wait timers, branch restrictions",
     "Custom deployment protection (ServiceNow, Datadog, etc.)"),
    ("Supply chain security",
     "Dependabot + CodeQL + Secret Scanning + Push Protection",
     "SLSA provenance attestations for build artifacts"),
    ("Fork PR isolation",
     "Secrets blocked from fork PRs; read-only GITHUB_TOKEN",
     "pull_request_target for controlled fork access (use with care)"),
]

y_start_sec = 1100000
y_step_sec = 780000
for i, (title, d1, d2) in enumerate(sec_items):
    add_bullet_item(slide_d, title, d1, d2,
                    left=450000, top=y_start_sec + i * y_step_sec,
                    bullet_shape=MSO_SHAPE.OVAL,
                    bullet_color=RGBColor(0x10, 0x7C, 0x10))

# Right: OIDC code example
code_oidc = """# OIDC — secretless Azure auth
permissions:
  id-token: write
  contents: read

steps:
  - uses: azure/login@v2
    with:
      client-id: ${{ vars.AZURE_CLIENT_ID }}
      tenant-id: ${{ vars.AZURE_TENANT_ID }}
      subscription-id: ${{ vars.AZURE_SUBSCRIPTION_ID }}
      # No client-secret needed!

# Environment with approval gate
jobs:
  deploy:
    environment:
      name: production    # requires reviewer approval
      url: https://myapp.azurewebsites.net
    steps:
      - uses: azure/webapps-deploy@v3

# Concurrency control
concurrency:
  group: deploy-${{ github.ref }}
  cancel-in-progress: false  # don't cancel deployments"""

add_code_block(slide_d, code_oidc,
               left=Emu(6200000), top=Emu(1100000),
               width=Emu(5600000), height=Emu(5400000),
               font_size=Pt(8))

ann_sec = (Emu(6200000), Emu(1050000), "Security patterns in practice")
txb = slide_d.shapes.add_textbox(ann_sec[0], ann_sec[1], Emu(4000000), Emu(200000))
p = txb.text_frame.paragraphs[0]
p.text = ann_sec[2]
run = p.runs[0]
run.font.size = Pt(10)
run.font.bold = True
run.font.color.rgb = RGBColor(0x10, 0x7C, 0x10)

# ═══════════════════════════════════════════
# SLIDE E: Azure DevOps vs GitHub Actions — Comparison
# ═══════════════════════════════════════════
slide_e = prs.slides.add_slide(blank_layout)
new_slides.append(slide_e)

add_left_bar(slide_e)
add_title(slide_e, "Azure DevOps Pipelines vs GitHub Actions", left=430225, top=100000)
add_title_underline(slide_e, left=512670, top=650000, width=10091420, height=36576)

# Comparison table
headers = ["Capability", "Azure DevOps Pipelines", "GitHub Actions"]
rows = [
    ["Pipeline hierarchy", "Stages → Jobs → Steps", "Jobs → Steps (no stages)"],
    ["Template reuse", "Steps/Jobs/Stages templates,\nextends keyword", "Reusable Workflows (jobs),\nComposite Actions (steps)"],
    ["Template enforcement", "✅ extends + required checks", "❌ No equivalent\n(org action restrictions only)"],
    ["Parameter types", "string, bool, number, object,\nstep/jobList, stageList", "string, boolean, number only\n(JSON workaround for objects)"],
    ["Dynamic generation", "✅ ${{ each }} over lists", "⚠️ Matrix + fromJSON()"],
    ["Expression phases", "Compile-time ${{ }}\n+ Runtime $[ ]", "Single runtime ${{ }}"],
    ["Secretless cloud auth", "⚠️ Workload Identity (preview)", "✅ OIDC federated credentials (GA)"],
    ["Secret masking", "Manual (issecret=true)", "Automatic (all secrets masked)"],
    ["Built-in SAST/SCA", "❌ Third-party required", "✅ CodeQL, Dependabot, Secret Scanning"],
    ["Marketplace ecosystem", "Azure DevOps Marketplace", "20,000+ GitHub Actions"],
    ["Free tier (private)", "1 parallel job, 1,800 min/mo\n~$40/mo per extra job", "2,000 min/mo (Linux)\n$0.008/min overage"],
    ["Self-hosted cost", "$15/mo per parallel slot", "Free & unlimited"],
    ["Max job timeout", "60 min free / 360 min paid", "360 min hosted / 35 days self-hosted"],
    ["Runner environments", "Linux, Windows, macOS\nNo GPU, No ARM64", "Linux, Windows, macOS\n+ ARM64, GPU (Enterprise)"],
    ["Scheduled triggers", "✅ Any branch", "❌ Default branch only"],
    ["Event types", "CI/CD + pipeline completion", "25+ (push, PR, issue, release, etc.)"],
]

add_table_slide(slide_e, headers, rows,
                left=Emu(300000), top=Emu(730000),
                width=Emu(11600000))


# ═══════════════════════════════════════════
# REORDER SLIDES: move new slides into position 14-18
# ═══════════════════════════════════════════

# After our deletions and additions, the slide list in XML is:
# Original slides 1-13, then 18-28 (renumbered to 14-24), then our new 5 slides at end (25-29)
# We want to insert our new slides at position 14 (after slide 13)

# Get the current slide ID list
slide_list = prs.slides._sldIdLst
all_sldIds = list(slide_list)

# The new slides are the last 5 entries
num_new = len(new_slides)
new_sldIds = all_sldIds[-num_new:]
old_sldIds = all_sldIds[:-num_new]

# Insert position: after index 12 (0-based), i.e., after slide 13
insert_pos = 13

# Build the new order
reordered = old_sldIds[:insert_pos] + new_sldIds + old_sldIds[insert_pos:]

# Clear and re-add in order
for sldId in list(slide_list):
    slide_list.remove(sldId)
for sldId in reordered:
    slide_list.append(sldId)

# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
prs.save(DST)
print(f"Saved to: {DST}")
print(f"Total slides: {len(list(prs.slides._sldIdLst))}")
