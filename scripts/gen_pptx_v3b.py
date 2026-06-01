"""
Generate v3 PowerPoint: Replace slides 14-17 content with GitHub Actions slides,
insert comparison slide after slide 17.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy
from lxml import etree
import shutil

SRC = r'C:\Users\darioa\Microsoft\Microsoft (Internal)-12 - Projects - PresentazioneDevops - Documents\AI Lab for Startup - GitHub + Azure - CI CD - v2a.pptx'
DST = r'C:\Users\darioa\Microsoft\Microsoft (Internal)-12 - Projects - PresentazioneDevops - Documents\AI Lab for Startup - GitHub + Azure - CI CD - v3.pptx'

# If DST is locked, use a temp path
import tempfile
try:
    shutil.copy2(SRC, DST)
    dst_path = DST
except PermissionError:
    dst_path = DST.replace('.pptx', '_new.pptx')
    shutil.copy2(SRC, dst_path)
    print(f"NOTE: Original file locked, writing to: {dst_path}")

prs = Presentation(dst_path)
blank_layout = prs.slide_layouts[6]  # Blank layout

# ═══════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════

def clear_slide(slide):
    """Remove all shapes from a slide."""
    sp_tree = slide.shapes._spTree
    for shape_elem in list(sp_tree):
        tag = etree.QName(shape_elem.tag).localname
        if tag in ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'):
            sp_tree.remove(shape_elem)

def add_left_bar(slide, width=164592):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1B, 0x1B, 0x1B)
    shape.line.fill.background()

def add_title(slide, text, left=430225, top=178385, width=10972800, height=646331):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

def add_title_underline(slide, left=512670, top=818464, width=10091420, height=45719):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x00, 0x78, 0xD4)
    shape.line.fill.background()

def add_section_header(slide, text, left, top, width=2028440, height=430887):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

def add_vertical_divider(slide, left=5989320, top=1097280, width=18288, height=5486400):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    shape.line.fill.background()

def add_bullet_item(slide, title_text, desc1, desc2, left, top,
                    bullet_shape=MSO_SHAPE.OVAL, bullet_color=RGBColor(0x00, 0x78, 0xD4),
                    title_size=Pt(13), desc_size=Pt(10)):
    # Bullet marker
    bshape = slide.shapes.add_shape(bullet_shape, left, top + 91440, 109728, 109728)
    bshape.fill.solid()
    bshape.fill.fore_color.rgb = bullet_color
    bshape.line.fill.background()
    
    # Text block
    txBox = slide.shapes.add_textbox(left + 182880, top, 5029200, 640080)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.runs[0].font.size = title_size
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    
    p2 = tf.add_paragraph()
    p2.text = desc1
    p2.runs[0].font.size = desc_size
    p2.runs[0].font.color.rgb = RGBColor(0x40, 0x40, 0x40)
    
    p3 = tf.add_paragraph()
    p3.text = desc2
    p3.runs[0].font.size = desc_size
    p3.runs[0].font.color.rgb = RGBColor(0x40, 0x40, 0x40)

def add_code_block(slide, code_text, left, top, width, height, font_size=Pt(8)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x1E)
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)
    
    for i, line in enumerate(code_text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = font_size
            run.font.name = 'Consolas'
            run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)

def add_annotation(slide, x, y, text, color=RGBColor(0x00, 0x78, 0xD4)):
    txb = slide.shapes.add_textbox(x, y, Emu(4000000), Emu(200000))
    p = txb.text_frame.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = color


# ═══════════════════════════════════════════
# SLIDE 14: GitHub Actions — At a Glance
# (replaces "CI/CD in GitHub" section header)
# ═══════════════════════════════════════════
slide_14 = prs.slides[13]
clear_slide(slide_14)
add_left_bar(slide_14)
add_title(slide_14, "GitHub Actions — At a Glance", left=429768, top=182880)
add_title_underline(slide_14, left=522706, top=797177, width=7262152, height=36576)
add_vertical_divider(slide_14)

add_section_header(slide_14, "Key Capabilities", left=548640, top=1051560)

cap_items = [
    ("Event-driven CI/CD",
     "25+ trigger types: push, PR, schedule, dispatch, release…",
     "Native GitHub integration — status checks, PR comments"),
    ("Reusable workflows + composite actions",
     "Modular: reusable workflows (jobs) + composite actions (steps)",
     "Cross-repo workflow calls from centralized template repos"),
    ("Service containers (sidecars)",
     "Spin up Postgres, Redis, etc. as sidecar containers per job",
     "Native integration testing — no external infra needed"),
    ("OIDC — secretless cloud auth",
     "Federated credentials for Azure, AWS, GCP — no stored secrets",
     "Token claims scoped by repo, branch, environment"),
    ("20,000+ marketplace actions",
     "Rich ecosystem: setup-dotnet, azure/login, docker/build-push…",
     "Composite, JavaScript, or Docker container action types"),
    ("Self-hosted runners (free & unlimited)",
     "No per-agent charge, no minute billing on your infrastructure",
     "Actions Runner Controller (ARC) for K8s auto-scaling"),
    ("Built-in supply chain security",
     "Dependabot, CodeQL SAST, secret scanning, push protection",
     "SLSA provenance attestations + dependency review in PRs"),
    ("GitHub Pages & Packages",
     "Built-in artifact hosting (npm, NuGet, Docker, Maven)",
     "GitHub Pages deployment with environment gates"),
]

y0 = 1490472
ys = 640080
for i, (t, d1, d2) in enumerate(cap_items):
    add_bullet_item(slide_14, t, d1, d2, left=594360, top=y0 + i * ys,
                    bullet_color=RGBColor(0x00, 0x78, 0xD4))

add_section_header(slide_14, "Key Limitations", left=6263640, top=1051560)

lim_items = [
    ("Free tier: 2,000 min/month (private repos)",
     "Linux ×1, Windows ×2, macOS ×10 minute multipliers",
     "Public repos: unlimited free minutes"),
    ("No stage concept",
     "Flat job graph only — stages simulated via needs: deps",
     "No visual stage grouping in the GitHub UI"),
    ("Reusable workflow depth: max 4 levels",
     "Cannot nest reusable workflows deeper than 4 calls",
     "Reusable workflows produce jobs, not individual steps"),
    ("Limited input types",
     "Only string, boolean, number — no object/array parameters",
     "Workaround: JSON-serialize via fromJSON()"),
    ("Job timeout: 6 hrs hosted, 35 days self-hosted",
     "GitHub-hosted max 360 min per job — no override available",
     "Self-hosted: up to 35 days per job"),
    ("Runner environments",
     "GitHub-hosted: Linux (x64/ARM64), Win (x64), macOS (x64/ARM64)",
     "Larger runners + GPU require Team/Enterprise plan"),
    ("Scheduled workflows: default branch only",
     "Cron triggers only run on the default branch",
     "Disabled after 60 days of repo inactivity"),
    ("No enforced template inheritance",
     "No extends keyword — cannot mandate approved templates",
     "Org-level action restrictions (allowlist) closest control"),
]

for i, (t, d1, d2) in enumerate(lim_items):
    add_bullet_item(slide_14, t, d1, d2, left=6309359, top=y0 + i * ys,
                    bullet_shape=MSO_SHAPE.DIAMOND,
                    bullet_color=RGBColor(0xD8, 0x3B, 0x01))


# ═══════════════════════════════════════════
# SLIDE 15: Structure of a GitHub Actions Workflow
# ═══════════════════════════════════════════
slide_15 = prs.slides[14]
clear_slide(slide_15)
add_left_bar(slide_15)
add_title(slide_15, "Structure of a GitHub Actions Workflow", left=430225, top=178385)
add_title_underline(slide_15)

code_left = """name: CI/CD Pipeline
on:
  push:
    branches: [main]
    paths: ['src/**']
  pull_request:
    branches: [main]
  workflow_dispatch:
    inputs:
      environment:
        type: choice
        options: [dev, staging, prod]

permissions:
  contents: read
  id-token: write

concurrency:
  group: ${{ github.workflow }}-${{ github.ref }}
  cancel-in-progress: true

env:
  DOTNET_VERSION: '9.0.x'"""

add_code_block(slide_15, code_left,
               left=Emu(300000), top=Emu(1000000),
               width=Emu(5400000), height=Emu(5500000))

code_right = """jobs:
  build:
    runs-on: ubuntu-latest
    outputs:
      version: ${{ steps.ver.outputs.version }}
    steps:
      - uses: actions/checkout@v4
      - uses: actions/setup-dotnet@v4
        with:
          dotnet-version: ${{ env.DOTNET_VERSION }}
      - run: dotnet build -c Release
      - id: ver
        run: echo "version=1.0.${{ github.run_number }}"
              >> $GITHUB_OUTPUT
      - uses: actions/upload-artifact@v4
        with:
          name: build-output
          path: ./publish/

  deploy:
    needs: build
    runs-on: ubuntu-latest
    environment:
      name: production
      url: https://myapp.azurewebsites.net
    steps:
      - uses: azure/login@v2
        with:
          client-id: ${{ vars.AZURE_CLIENT_ID }}
          tenant-id: ${{ vars.AZURE_TENANT_ID }}
          subscription-id: ${{ vars.AZURE_SUBSCRIPTION_ID }}
      - uses: azure/webapps-deploy@v3
        with:
          app-name: my-app"""

add_code_block(slide_15, code_right,
               left=Emu(5900000), top=Emu(1000000),
               width=Emu(5900000), height=Emu(5500000))

add_annotation(slide_15, Emu(300000), Emu(920000), "Triggers, permissions & concurrency")
add_annotation(slide_15, Emu(5900000), Emu(920000), "Jobs, steps & deployment with OIDC")


# ═══════════════════════════════════════════
# SLIDE 16: Reusable Workflows & Patterns
# ═══════════════════════════════════════════
slide_16 = prs.slides[15]
clear_slide(slide_16)
add_left_bar(slide_16)
add_title(slide_16, "GitHub Actions — Reusable Workflows & Patterns", left=430225, top=178385)
add_title_underline(slide_16)

code_reusable = """# Reusable workflow (.github/workflows/deploy.yml)
name: Deploy to Azure
on:
  workflow_call:
    inputs:
      environment:
        required: true
        type: string
      app-name:
        required: true
        type: string
    outputs:
      url:
        value: ${{ jobs.deploy.outputs.url }}

jobs:
  deploy:
    runs-on: ubuntu-latest
    environment:
      name: ${{ inputs.environment }}
    outputs:
      url: ${{ steps.deploy.outputs.webapp-url }}
    steps:
      - uses: azure/login@v2
        with:
          client-id: ${{ vars.AZURE_CLIENT_ID }}
          tenant-id: ${{ vars.AZURE_TENANT_ID }}
          subscription-id: ${{ vars.AZURE_SUBSCRIPTION_ID }}
      - uses: azure/webapps-deploy@v3
        id: deploy
        with:
          app-name: ${{ inputs.app-name }}"""

add_code_block(slide_16, code_reusable,
               left=Emu(300000), top=Emu(1000000),
               width=Emu(5400000), height=Emu(4600000))

code_caller = """# Caller workflow — multi-env deployment
jobs:
  deploy-staging:
    uses: ./.github/workflows/deploy.yml
    with:
      environment: staging
      app-name: my-app-staging
    secrets: inherit

  deploy-prod:
    needs: deploy-staging
    uses: ./.github/workflows/deploy.yml
    with:
      environment: production
      app-name: my-app-prod
    secrets: inherit

# Matrix strategy — cross-platform testing
  test:
    strategy:
      fail-fast: false
      matrix:
        os: [ubuntu-latest, windows-latest]
        dotnet: ['8.0.x', '9.0.x']
    runs-on: ${{ matrix.os }}
    steps:
      - uses: actions/setup-dotnet@v4
        with:
          dotnet-version: ${{ matrix.dotnet }}
      - run: dotnet test"""

add_code_block(slide_16, code_caller,
               left=Emu(5900000), top=Emu(1000000),
               width=Emu(5900000), height=Emu(4600000))

add_annotation(slide_16, Emu(300000), Emu(920000), "Reusable workflow definition (callee)")
add_annotation(slide_16, Emu(5900000), Emu(920000), "Caller workflow + matrix strategy")

# Key patterns summary
patterns = [
    "▸ Reusable workflows (workflow_call) — job-level reuse across repos    |    ▸ Composite actions — step-level reuse (local or cross-repo)",
    "▸ Matrix strategy — N parallel jobs from OS × version × config           |    ▸ Dynamic matrix — fromJSON() for runtime-computed job grids",
    "▸ OIDC auth — secretless Azure/AWS/GCP via federated credentials     |    ▸ secrets: inherit — pass all caller secrets implicitly",
]
txBox = slide_16.shapes.add_textbox(Emu(300000), Emu(5750000), Emu(11500000), Emu(900000))
tf = txBox.text_frame
tf.word_wrap = True
for i, pt in enumerate(patterns):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = pt
    p.runs[0].font.size = Pt(9)
    p.runs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)


# ═══════════════════════════════════════════
# SLIDE 17: Security & Environments
# ═══════════════════════════════════════════
slide_17 = prs.slides[16]
clear_slide(slide_17)
add_left_bar(slide_17)
add_title(slide_17, "GitHub Actions — Security & Environments", left=430225, top=178385)
add_title_underline(slide_17)

sec_items = [
    ("OIDC Federated Credentials",
     "Secretless auth to Azure, AWS, GCP — no stored credentials",
     "Token scoped by repo + branch + environment claims"),
    ("GITHUB_TOKEN (auto-provisioned)",
     "Per-workflow token with configurable permission scopes",
     "Least-privilege: contents:read, id-token:write, etc."),
    ("Secrets at 3 levels",
     "Organization → Repository → Environment scoped secrets",
     "Masked in logs, blocked from fork PRs"),
    ("Environment protection rules",
     "Required reviewers, wait timers, branch restrictions",
     "Custom protection rules (ServiceNow, Datadog, etc.)"),
    ("Supply chain security",
     "Dependabot + CodeQL + Secret Scanning + Push Protection",
     "SLSA provenance attestations for build artifacts"),
    ("Fork PR isolation",
     "Secrets blocked from fork PRs; read-only GITHUB_TOKEN",
     "pull_request_target for controlled fork access"),
]

for i, (t, d1, d2) in enumerate(sec_items):
    add_bullet_item(slide_17, t, d1, d2,
                    left=450000, top=1100000 + i * 780000,
                    bullet_color=RGBColor(0x10, 0x7C, 0x10))

code_sec = """# OIDC — secretless Azure authentication
permissions:
  id-token: write
  contents: read

steps:
  - uses: azure/login@v2
    with:
      client-id: ${{ vars.AZURE_CLIENT_ID }}
      tenant-id: ${{ vars.AZURE_TENANT_ID }}
      subscription-id: ${{ vars.AZURE_SUBSCRIPTION_ID }}
      # No client-secret needed!

# Environment with approval gate
jobs:
  deploy:
    environment:
      name: production  # requires approval
      url: https://myapp.azurewebsites.net
    steps:
      - uses: azure/webapps-deploy@v3

# Concurrency — prevent overlapping deploys
concurrency:
  group: deploy-${{ github.ref }}
  cancel-in-progress: false"""

add_code_block(slide_17, code_sec,
               left=Emu(6200000), top=Emu(1100000),
               width=Emu(5600000), height=Emu(5400000))

add_annotation(slide_17, Emu(6200000), Emu(1050000), "Security patterns in practice",
               color=RGBColor(0x10, 0x7C, 0x10))


# ═══════════════════════════════════════════
# NEW SLIDE 18: Comparison — Azure DevOps vs GitHub Actions
# Insert after slide 17, before AVM section
# ═══════════════════════════════════════════
slide_cmp = prs.slides.add_slide(blank_layout)

add_left_bar(slide_cmp)
add_title(slide_cmp, "Azure DevOps vs GitHub Actions — Comparison", left=430225, top=80000)
add_title_underline(slide_cmp, left=512670, top=620000, width=10091420, height=36576)

# Build the comparison table
headers = ["Capability", "Azure DevOps Pipelines", "GitHub Actions"]
rows = [
    ["Pipeline hierarchy", "Stages → Jobs → Steps", "Jobs → Steps (no stages)"],
    ["Template reuse", "Steps/Jobs/Stages templates\nextends keyword", "Reusable Workflows (jobs)\nComposite Actions (steps)"],
    ["Template enforcement", "✅ extends + required checks", "❌ No equivalent (org restrictions)"],
    ["Parameter types", "string, bool, number, object,\nstepList, stageList", "string, boolean, number only\n(JSON workaround for objects)"],
    ["Dynamic generation", "✅ ${{ each }} over lists", "⚠️ Matrix + fromJSON()"],    ("Service containers", "❌ Not available natively", "✅ Native sidecar containers per job"),    ["Expression phases", "Compile-time ${{ }}\n+ Runtime $[ ]", "Single runtime ${{ }}"],
    ["Secretless cloud auth", "⚠️ Workload Identity (preview)", "✅ OIDC (GA)"],
    ["Built-in SAST/SCA", "❌ Third-party required", "✅ CodeQL, Dependabot, etc."],
    ["Marketplace", "Azure DevOps Marketplace", "20,000+ GitHub Actions"],
    ["Free (private repos)", "1 job, 1,800 min/mo\n~$40/mo per extra job", "2,000 min/mo (Linux)\n$0.008/min overage"],
    ["Self-hosted cost", "$15/mo per parallel slot", "Free & unlimited"],
    ["Max job timeout", "60 min free / 360 min paid", "360 min hosted / 35d self-hosted"],
    ["Runners", "Linux, Win, macOS\nNo GPU, No ARM64", "Linux, Win, macOS\n+ ARM64, GPU (Enterprise)"],
    ["Scheduled triggers", "✅ Any branch", "❌ Default branch only"],
    ["Event types", "CI/CD + pipeline completion", "25+ (push, PR, issue, release…)"],
]

num_rows = len(rows) + 1
num_cols = 3
tbl_left = Emu(300000)
tbl_top = Emu(700000)
tbl_width = Emu(11600000)
row_h = Emu(300000)
tbl_h = row_h * num_rows

tbl_shape = slide_cmp.shapes.add_table(num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_h)
table = tbl_shape.table

# Set column widths
table.columns[0].width = Emu(2500000)
table.columns[1].width = Emu(4550000)
table.columns[2].width = Emu(4550000)

# Header row
for ci, h in enumerate(headers):
    cell = table.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x00, 0x78, 0xD4)

# Data rows
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = table.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        if ri % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ═══════════════════════════════════════════
# REORDER: Move comparison slide to position 18 (after slide 17)
# ═══════════════════════════════════════════
slide_list = prs.slides._sldIdLst
all_ids = list(slide_list)

# The new comparison slide is the last one
cmp_id = all_ids[-1]
all_ids.pop()

# Insert at position 17 (0-based index 17 = slide 18)
all_ids.insert(17, cmp_id)

# Rebuild slide list
for sid in list(slide_list):
    slide_list.remove(sid)
for sid in all_ids:
    slide_list.append(sid)

# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
prs.save(dst_path)
print(f"Saved to: {dst_path}")
print(f"Total slides: {len(list(prs.slides._sldIdLst))}")

# Verify
prs2 = Presentation(dst_path)
for idx, slide in enumerate(prs2.slides):
    title = ''
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t and len(t) > 5:
                    title = t[:80]
                    break
        if title:
            break
    print(f'  Slide {idx+1}: {title}')
