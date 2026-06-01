from pptx import Presentation
from pptx.util import Inches, Pt, Emu

pptx_path = r'C:\Users\darioa\Microsoft\Microsoft (Internal)-12 - Projects - PresentazioneDevops - Documents\AI Lab for Startup - GitHub + Azure - CI CD - v2a.pptx'
prs = Presentation(pptx_path)

print(f'Slide width: {prs.slide_width}, height: {prs.slide_height}')
print(f'Total slides: {len(prs.slides)}')
print(f'Slide layouts available: {len(prs.slide_layouts)}')
for i, layout in enumerate(prs.slide_layouts):
    print(f'  Layout {i}: {layout.name}')
print()

for idx, slide in enumerate(prs.slides):
    print(f'=== Slide {idx+1} ===')
    print(f'  Layout: {slide.slide_layout.name}')
    for shape in slide.shapes:
        print(f'  Shape: {shape.shape_type}, name={shape.name}, pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})')
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    font_info = ''
                    if para.runs:
                        r = para.runs[0]
                        font_info = f' [font={r.font.name}, size={r.font.size}, bold={r.font.bold}]'
                    print(f'    "{text[:150]}"{font_info}')
        if shape.has_table:
            table = shape.table
            print(f'    TABLE: {len(table.rows)} rows x {len(table.columns)} cols')
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip()[:60] for cell in row.cells]
                print(f'      Row {row_idx}: {cells}')
    print()
